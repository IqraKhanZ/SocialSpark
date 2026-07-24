"""
SocialSpark Premium Pitch Deck — v3
Fixes: images now visible, emojis replaced with geometric icon shapes,
Thank You layout corrected, content sharpened for business presentation.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ───────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
ORANGE     = RGBColor(0xF6, 0x3B, 0x05)
CREAM      = RGBColor(0xFF, 0xFA, 0xE5)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x14, 0x14, 0x14)
MID_GRAY   = RGBColor(0x2A, 0x2A, 0x2A)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)

FONT_DISPLAY = "Anton"
FONT_BODY    = "Open Sans"

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR  = os.path.join(BASE_DIR, "images")
OUTPUT   = os.path.join(BASE_DIR, "SocialSpark_PitchDeck_v3.pptx")

def img(name):
    return os.path.join(IMG_DIR, name + ".jpg")

# ── Base helpers ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(20)
    prs.slide_height = Inches(11.25)
    return prs

def set_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def add_rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        _solid(s, fill)
    else:
        s.fill.background()
    no_line(s)
    return s

def add_rrect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        _solid(s, fill)
    else:
        s.fill.background()
    no_line(s)
    return s

def add_pic(slide, path, l, t, w, h):
    """Add picture only if file exists."""
    if os.path.exists(path):
        pic = slide.shapes.add_picture(
            path, Inches(l), Inches(t), Inches(w), Inches(h))
        return pic
    return None

def add_pic_with_alpha_overlay(slide, path, l, t, w, h, overlay_color, alpha_pct=50):
    """Add picture, then add a semi-transparent overlay via XML alpha."""
    pic = add_pic(slide, path, l, t, w, h)
    # Draw overlay shape on top
    ov = add_rect(slide, l, t, w, h, fill=overlay_color)
    # Set alpha via XML
    spPr = ov._element.find(qn('p:spPr'))
    if spPr is not None:
        solidFill = spPr.find('.//' + qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', str(int(alpha_pct * 1000)))
    return pic, ov

def add_text(slide, text, l, t, w, h,
             font=None, size=Pt(18), bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    font = font or FONT_BODY
    txb = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap   = wrap
    tf.margin_left = tf.margin_right  = Inches(0.0)
    tf.margin_top  = tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text        = text
    r.font.name   = font
    r.font.size   = size
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb, tf

def add_rule(slide, l, t, w, thickness=0.04, color=ORANGE):
    add_rect(slide, l, t, w, thickness, fill=color)

def add_label_tag(slide, text, l, t, w=3.8, h=0.42):
    """Orange tag with uppercase label text."""
    add_rect(slide, l, t, w, h, fill=ORANGE)
    add_text(slide, text.upper(), l + 0.15, t + 0.07, w - 0.3, h - 0.14,
             size=Pt(13), bold=True, color=WHITE)

def add_footer(slide, num, total=12):
    add_rect(slide, 0, 10.85, 20, 0.4, fill=DARK_GRAY)
    add_text(slide,
             "SocialSpark  —  Building Social Fitness, One Connection at a Time.",
             0.4, 10.9, 15, 0.3, size=Pt(11), color=LIGHT_GRAY)
    add_text(slide, f"{num:02d} / {total}",
             18.2, 10.9, 1.6, 0.3,
             font=FONT_DISPLAY, size=Pt(14), color=ORANGE,
             align=PP_ALIGN.RIGHT)

# ── Icon shape helper (replaces emojis) ─────────────────────────────────
def add_icon_shape(slide, symbol, l, t, size_in=0.7,
                   bg_color=ORANGE, text_color=WHITE, text_size=Pt(22)):
    """
    Draws a small filled square with a bold character/symbol as a professional icon.
    'symbol' should be a short string: '01', 'XP', a Unicode geometric char, etc.
    """
    sq = add_rect(slide, l, t, size_in, size_in, fill=bg_color)
    add_text(slide, symbol,
             l, t + (size_in - size_in * 0.55) / 2,
             size_in, size_in * 0.6,
             font=FONT_DISPLAY, size=text_size,
             bold=False, color=text_color, align=PP_ALIGN.CENTER)

def add_icon_circle(slide, symbol, l, t, size_in=0.7,
                    bg_color=ORANGE, text_color=WHITE, text_size=Pt(20)):
    """Oval-based icon for variety."""
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(l), Inches(t), Inches(size_in), Inches(size_in))
    _solid(circ, bg_color)
    no_line(circ)
    add_text(slide, symbol,
             l, t + (size_in * 0.18),
             size_in, size_in * 0.65,
             font=FONT_DISPLAY, size=text_size,
             bold=False, color=text_color, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────
# SLIDES
# ─────────────────────────────────────────────────────────────────────────

def s01_cover(prs):
    """SLIDE 1 — Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    # Full right-side image — NO overlay so image is visible
    add_pic(slide, img("community_connection"), 10.0, 0, 10.0, 11.25)

    # Left black content panel
    add_rect(slide, 0, 0, 10.2, 11.25, fill=BLACK)
    # Left thin orange stripe accent
    add_rect(slide, 0, 0, 0.45, 11.25, fill=ORANGE)
    # Orange diagonal accent bar at bottom of left panel
    add_rect(slide, 0, 9.5, 10.2, 0.07, fill=ORANGE)

    # Big headline
    add_text(slide, "SOCIAL", 0.8, 0.6, 9.2, 2.8,
             font=FONT_DISPLAY, size=Pt(130), color=WHITE)
    add_text(slide, "SPARK", 0.8, 3.3, 9.2, 2.8,
             font=FONT_DISPLAY, size=Pt(130), color=ORANGE)

    add_rule(slide, 0.8, 6.4, 5.5)
    add_text(slide,
             "Building Social Fitness,\nOne Connection at a Time.",
             0.8, 6.6, 9.0, 1.5,
             size=Pt(22), color=CREAM)
    add_text(slide, "Presented by: Team SocialSpark",
             0.8, 8.4, 9.0, 0.5, size=Pt(16), color=LIGHT_GRAY)
    add_text(slide, "Entrepreneurship Cell  —  Shark Tank Competition",
             0.8, 8.95, 9.0, 0.5, size=Pt(14), color=ORANGE)

    add_footer(slide, 1)


def s02_problem(prs):
    """SLIDE 2 — The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    # Left image: no overlay — fully visible
    add_pic(slide, img("loneliness_problem"), 0, 0, 10.5, 11.25)

    # Gradient-like dark strip on far left so text can sit on top of image
    # Use a thin black-to-transparent approach: just place text on a semi-dark bar at bottom
    add_rect(slide, 0, 7.0, 10.5, 4.25, fill=BLACK)

    # Right content panel
    add_rect(slide, 10.5, 0, 9.5, 11.25, fill=DARK_GRAY)
    add_rect(slide, 10.5, 0, 0.07, 11.25, fill=ORANGE)

    # Over-image quote (bottom of image)
    add_text(slide,
             '"We don\'t lack people around us.\nWe lack meaningful connections."',
             0.5, 7.2, 9.5, 2.3,
             size=Pt(22), italic=True, color=WHITE)

    # Right panel
    add_label_tag(slide, "The Problem", 11.0, 0.65, 3.8)
    add_text(slide, "LONELINESS\nIS AN ECONOMIC\nCRISIS.", 11.0, 1.35, 8.2, 4.5,
             font=FONT_DISPLAY, size=Pt(52), color=WHITE)
    add_rule(slide, 11.0, 6.1, 8.2)

    # Stats using icon shapes instead of emoji
    stats = [
        ("1:4",  "Adults worldwide experience chronic loneliness."),
        ("3x",   "Higher risk of depression, anxiety & burnout."),
        ("16-60", "The full age range impacted in India alone."),
    ]
    y = 6.4
    for sym, desc in stats:
        add_icon_shape(slide, sym, 11.0, y, 0.8, bg_color=ORANGE, text_size=Pt(14))
        add_text(slide, desc, 12.1, y + 0.12, 7.0, 0.58,
                 size=Pt(15), color=CREAM)
        y += 1.0

    add_footer(slide, 2)


def s03_gap(prs):
    """SLIDE 3 — Why Current Solutions Fail"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    add_rect(slide, 0, 0, 20, 0.07, fill=ORANGE)
    add_label_tag(slide, "The Gap", 0.8, 0.35, 2.8)
    add_text(slide, "COMMUNICATION\nVS CONNECTION.", 0.8, 0.95, 13, 2.6,
             font=FONT_DISPLAY, size=Pt(68), color=WHITE)
    add_text(slide,
             "Apps help people talk. SocialSpark helps people belong.",
             0.8, 3.7, 13, 0.65,
             size=Pt(20), bold=True, color=ORANGE)
    add_rule(slide, 0.8, 4.55, 18.5)

    # LEFT column header
    add_rect(slide, 0.8, 4.85, 8.8, 0.55, fill=MID_GRAY)
    add_text(slide, "Existing Platforms", 1.0, 4.87, 8.4, 0.5,
             font=FONT_DISPLAY, size=Pt(22), color=LIGHT_GRAY)
    rows_l = [
        "Chat-first — conversation without real action",
        "Screen-addictive — passive, online-only loops",
        "Anxiety-inducing — vanity metrics drive comparison",
        "No offline accountability or habit formation",
        "Connections are weak — no shared experience",
    ]
    y = 5.55
    for row in rows_l:
        add_text(slide, "—  " + row, 1.0, y, 8.4, 0.62, size=Pt(15), color=CREAM)
        y += 0.72

    # RIGHT column header
    add_rect(slide, 10.1, 4.85, 9.0, 0.55, fill=ORANGE)
    add_text(slide, "SocialSpark", 10.35, 4.87, 8.4, 0.5,
             font=FONT_DISPLAY, size=Pt(22), color=BLACK)
    rows_r = [
        "Experience-first — real-world offline events",
        "QR-verified attendance — no screen addiction",
        "Gamified XP — positive reinforcement system",
        "Streaks & missions — builds long-term habits",
        "Bonds through shared activities, not chatting",
    ]
    y = 5.55
    for row in rows_r:
        add_text(slide, "+  " + row, 10.35, y, 8.4, 0.62,
                 size=Pt(15), bold=True, color=WHITE)
        y += 0.72

    add_footer(slide, 3)


def s04_solution(prs):
    """SLIDE 4 — Introducing SocialSpark"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    # Right panel: app mockup — visible, no solid overlay
    add_pic(slide, img("app_mockup_screen"), 11.2, 0.4, 8.3, 10.45)

    # Left content zone
    add_rect(slide, 0, 0, 11.4, 11.25, fill=BLACK)
    add_rect(slide, 0, 0, 0.45, 11.25, fill=ORANGE)

    add_label_tag(slide, "Our Solution", 0.8, 0.7, 3.5)
    add_text(slide, "INDIA'S FIRST\nSOCIAL FITNESS\nPLATFORM.", 0.8, 1.4, 10.4, 4.2,
             font=FONT_DISPLAY, size=Pt(60), color=WHITE)

    add_rule(slide, 0.8, 5.85, 10.4)

    add_text(slide, "Instead of asking who you want to talk to,",
             0.8, 6.1, 10.0, 0.55,
             size=Pt(16), italic=True, color=LIGHT_GRAY)
    add_text(slide, 'We ask: "What do you want to experience today?"',
             0.8, 6.75, 10.0, 0.65,
             size=Pt(18), bold=True, color=ORANGE)

    features = [
        ("HC", "Hobby Clubs", "Photography, Chess, Music, Cooking"),
        ("CE", "Community Events", "Sports, Tournaments, Fitness Drives"),
        ("VD", "Volunteer Drives", "NGO-verified, XP-rewarded"),
        ("LW", "Local Workshops", "Coding, Pottery, Dance, Language"),
    ]
    y = 7.6
    for sym, title, sub in features:
        add_icon_shape(slide, sym, 0.8, y, 0.55, text_size=Pt(12))
        add_text(slide, title, 1.55, y, 4.0, 0.3,
                 size=Pt(15), bold=True, color=WHITE)
        add_text(slide, sub, 1.55, y + 0.32, 9.0, 0.3,
                 size=Pt(12), color=LIGHT_GRAY)
        y += 0.74

    add_footer(slide, 4)


def s05_howit(prs):
    """SLIDE 5 — How It Works"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    add_rect(slide, 0, 0, 20, 0.07, fill=ORANGE)
    add_label_tag(slide, "The Product", 0.8, 0.35, 3.0)
    add_text(slide, "HOW IT WORKS.", 0.8, 0.95, 15, 1.8,
             font=FONT_DISPLAY, size=Pt(82), color=WHITE)
    add_text(slide, "Four steps. Fully verified. Beautifully simple.",
             0.8, 2.85, 12, 0.6, size=Pt(20), color=ORANGE)
    add_rule(slide, 0.8, 3.65, 18.5)

    steps = [
        ("01", "SIGN UP\n& SELECT",
         "Pick your world — Photography, Football, Yoga, Cooking, Music, "
         "Coding, Reading, or Volunteering."),
        ("02", "AI\nMATCHES",
         "Our algorithm surfaces the right verified, nearby events based on "
         "your interest profile and location."),
        ("03", "ATTEND\n& SCAN",
         "Show up. Scan the host QR code. Attendance is instantly "
         "verified — no self-reporting, no cheating."),
        ("04", "EARN XP\n& WIN",
         "Collect Social XP. Hit streaks. Climb leaderboards. Unlock "
         "real rewards — coffee, cinema, gyms, classes."),
    ]
    card_w = 4.5
    x = 0.8
    for num, title, desc in steps:
        add_rect(slide, x, 4.0, card_w - 0.2, 5.85, fill=MID_GRAY)
        add_rect(slide, x, 4.0, card_w - 0.2, 0.06, fill=ORANGE)
        add_text(slide, num, x + 0.2, 4.15, 1.8, 0.95,
                 font=FONT_DISPLAY, size=Pt(44), color=ORANGE)
        add_text(slide, title, x + 0.2, 5.15, card_w - 0.5, 1.1,
                 font=FONT_DISPLAY, size=Pt(24), color=WHITE)
        add_text(slide, desc, x + 0.2, 6.4, card_w - 0.5, 3.2,
                 size=Pt(14), color=LIGHT_GRAY)
        x += card_w

    add_footer(slide, 5)


def s06_gamify(prs):
    """SLIDE 6 — Social Fitness Engine"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    # Right image panel — NO opaque overlay so image shows
    add_pic(slide, img("gamification_xp"), 10.5, 0, 9.5, 11.25)
    # Only add a narrow dark edge strip on the LEFT side of the image
    # to create visual separation (not covering the image)
    add_rect(slide, 10.5, 0, 0.08, 11.25, fill=ORANGE)

    add_rect(slide, 0, 0, 10.42, 11.25, fill=BLACK)
    add_rect(slide, 0, 0, 0.45, 11.25, fill=ORANGE)

    add_label_tag(slide, "Social Fitness Engine", 0.8, 0.55, 4.8)
    add_text(slide, "MAKING SOCIAL\nINTERACTION\nA HABIT.", 0.8, 1.2, 9.5, 3.6,
             font=FONT_DISPLAY, size=Pt(64), color=WHITE)

    add_rule(slide, 0.8, 5.0, 9.5)

    # XP table — icons as small labeled squares
    xp_rows = [
        ("ATT", "Attend any verified event",           "+50 XP"),
        ("VOL", "Volunteer at an NGO drive",            "+150 XP"),
        ("HST", "Host a community experience",         "+300 XP"),
        ("REF", "Refer and bring a new friend",         "+100 XP"),
        ("QST", "Complete the weekly social quest",    "+250 XP"),
    ]
    y = 5.3
    for sym, action, xp in xp_rows:
        add_rect(slide, 0.8, y, 9.5, 0.62, fill=MID_GRAY)
        add_icon_shape(slide, sym, 0.85, y + 0.06, 0.52, bg_color=ORANGE, text_size=Pt(11))
        add_text(slide, action, 1.55, y + 0.14, 6.5, 0.42,
                 size=Pt(15), color=CREAM)
        add_text(slide, xp, 8.1, y + 0.1, 2.1, 0.48,
                 font=FONT_DISPLAY, size=Pt(20), color=ORANGE,
                 align=PP_ALIGN.RIGHT)
        y += 0.74

    add_text(slide,
             "Exactly how RunSpark built exercise habits —\n"
             "SocialSpark builds lifelong social fitness habits.",
             0.8, 9.8, 9.5, 1.0,
             size=Pt(14), italic=True, color=LIGHT_GRAY)

    add_footer(slide, 6)


def s07_business(prs):
    """SLIDE 7 — Business Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    add_rect(slide, 0, 0, 20, 0.07, fill=ORANGE)
    add_label_tag(slide, "Business Model", 0.8, 0.35, 3.8)
    add_text(slide, "FIVE REVENUE\nSTREAMS.", 0.8, 0.9, 12, 2.1,
             font=FONT_DISPLAY, size=Pt(74), color=WHITE)
    add_text(slide, "Diversified  —  Recurring  —  Scalable.",
             0.8, 3.05, 12, 0.6, size=Pt(20), color=ORANGE)
    add_rule(slide, 0.8, 3.8, 18.5)

    streams = [
        ("B2B", "B2B Hosting Fees",
         "Cafes, bookstores, and coworking spaces pay a listing and promotion fee "
         "to host branded experiences and drive real offline footfall — no billboard spend."),
        ("PRE", "Premium Membership",
         "Rs 199/month unlocks priority event bookings, exclusive workshops, a 2x XP "
         "multiplier, and advanced community networking tools."),
        ("COM", "Ticket Commission",
         "20% cut on all paid hobby workshops — pottery, dance, coding, music, "
         "photography, cooking. Every event is recurring revenue."),
        ("BRD", "Brand Partnerships",
         "Local and national brands sponsor XP leaderboards, seasonal campaigns, "
         "and community challenges instead of expensive influencer deals."),
        ("NGO", "NGO Subscriptions",
         "Monthly SaaS fee for NGOs to recruit, verify, and manage dedicated volunteer "
         "cohorts through our platform with mission tracking."),
    ]
    y = 4.15
    for sym, title, desc in streams:
        add_icon_shape(slide, sym, 0.8, y, 0.6, text_size=Pt(12))
        add_text(slide, title, 1.65, y, 4.5, 0.38,
                 size=Pt(16), bold=True, color=ORANGE)
        add_text(slide, desc, 1.65, y + 0.42, 17.5, 0.5,
                 size=Pt(13), color=CREAM)
        add_rule(slide, 0.8, y + 1.08, 18.5, thickness=0.018, color=MID_GRAY)
        y += 1.2

    add_footer(slide, 7)


def s08_scale(prs):
    """SLIDE 8 — Market & Scalability"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_GRAY)

    # Left map image — NO opaque overlay. Show it clean.
    add_pic(slide, img("india_scale_map"), 0, 0, 10.5, 11.25)

    # Right content panel
    add_rect(slide, 10.5, 0, 9.5, 11.25, fill=BLACK)
    add_rect(slide, 10.5, 0, 0.07, 11.25, fill=ORANGE)

    add_label_tag(slide, "Scalability Roadmap", 11.0, 0.55, 4.8)
    add_text(slide, "BUILT\nTO SCALE.", 11.0, 1.25, 8.5, 2.3,
             font=FONT_DISPLAY, size=Pt(72), color=WHITE)
    add_rule(slide, 11.0, 3.75, 8.5)

    phases = [
        ("01", "PILOT CITY",
         "20 cafes, 15 NGOs, 10 colleges, 100 monthly events."),
        ("02", "CAMPUS & CORPORATES",
         "Universities, hybrid workplaces, and coworking networks."),
        ("03", "METRO CITIES",
         "Bangalore, Mumbai, Delhi, Hyderabad, Pune."),
        ("04", "PAN INDIA",
         "Tier 1 and Tier 2 cities. National community scale."),
        ("05", "GLOBAL EXPANSION",
         "South-East Asia, Middle East, and international student hubs."),
    ]
    y = 4.0
    for num, phase, detail in phases:
        add_icon_shape(slide, num, 11.0, y, 0.62, text_size=Pt(14))
        add_text(slide, phase, 11.85, y + 0.02, 3.8, 0.42,
                 font=FONT_DISPLAY, size=Pt(18), color=ORANGE)
        add_text(slide, detail, 11.85, y + 0.46, 7.5, 0.42,
                 size=Pt(13), color=CREAM)
        add_rule(slide, 11.0, y + 1.08, 8.5, thickness=0.02, color=MID_GRAY)
        y += 1.22

    add_footer(slide, 8)


def s09_usp(prs):
    """SLIDE 9 — Why SocialSpark (USPs)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    add_rect(slide, 0, 0, 20, 0.07, fill=ORANGE)
    add_label_tag(slide, "Competitive Advantage", 0.8, 0.35, 4.5)
    add_text(slide, "FIVE REASONS\nWE WIN.", 0.8, 0.9, 12, 2.1,
             font=FONT_DISPLAY, size=Pt(68), color=WHITE)
    add_rule(slide, 0.8, 3.1, 18.5)

    usps = [
        ("SF",  "SOCIAL FITNESS",
         "We don't compete with Instagram.\nWe track and reward what Instagram ignores — real offline human connection."),
        ("OF",  "OFFLINE FIRST",
         "Technology is just the trigger.\nReal friendships are forged through shared experiences, not shared feeds."),
        ("XP",  "GAMIFIED HABITS",
         "XP, streaks, leaderboards, weekly quests.\nThe psychology that made fitness apps indispensable — applied to social life."),
        ("VF",  "VERIFIED & SAFE",
         "Verified profiles, verified hosts, verified venues.\nTrust is not an add-on. It is the foundation of every interaction."),
        ("EC",  "FULL ECOSYSTEM",
         "Users, businesses, NGOs, and brands.\nEvery stakeholder is rewarded, creating a self-sustaining value network."),
    ]
    card_w = 3.68
    x = 0.8
    for sym, title, desc in usps:
        add_rect(slide, x, 3.4, card_w - 0.1, 6.35, fill=MID_GRAY)
        add_rect(slide, x, 3.4, card_w - 0.1, 0.06, fill=ORANGE)
        add_icon_circle(slide, sym, x + 0.3, 3.65, 0.82,
                        bg_color=ORANGE, text_size=Pt(18))
        add_text(slide, title, x + 0.15, 4.65, card_w - 0.35, 0.75,
                 font=FONT_DISPLAY, size=Pt(21), color=WHITE)
        add_text(slide, desc, x + 0.15, 5.5, card_w - 0.35, 3.85,
                 size=Pt(13), color=LIGHT_GRAY)
        x += card_w

    add_text(slide, '"We don\'t build users. We build communities."',
             0.8, 9.82, 18.5, 0.75,
             size=Pt(18), italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

    add_footer(slide, 9)


def s10_whynow(prs):
    """SLIDE 10 — Why Now / Market Timing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    # Left dark statement panel
    add_rect(slide, 0, 0, 9.5, 11.25, fill=BLACK)
    add_rect(slide, 0, 0, 0.45, 11.25, fill=ORANGE)
    add_rect(slide, 9.5, 0, 0.07, 11.25, fill=ORANGE)

    add_label_tag(slide, "Market Timing", 0.8, 0.65, 3.5)
    add_text(slide, "RIGHT\nPROBLEM.\nRIGHT\nTIME.", 0.8, 1.35, 8.3, 5.2,
             font=FONT_DISPLAY, size=Pt(68), color=WHITE)
    add_rule(slide, 0.8, 6.8, 8.3)
    add_text(slide,
             '"People track sleep, calories, steps.\nWhy not their social wellbeing?"',
             0.8, 7.05, 8.3, 2.5,
             size=Pt(20), italic=True, color=CREAM)

    # Three right columns — each a distinct panel
    col_data = [
        ("MH",  "MENTAL HEALTH\nAWARENESS",
         "Loneliness has been declared a global health epidemic. WHO and governments "
         "worldwide are actively funding community health solutions.",
         DARK_GRAY),
        ("HW",  "HYBRID WORK\nIS THE NORM",
         "Remote and hybrid work has eliminated the daily spontaneous social contact "
         "of an office. The demand for real community has never been greater.",
         MID_GRAY),
        ("SF",  "SCREEN\nFATIGUE",
         "Gen-Z and Millennials are actively seeking digital detoxes. "
         "Real, tangible experiences are becoming the new social currency.",
         DARK_GRAY),
    ]
    col_x = [9.9, 13.3, 16.7]
    for i, (sym, title, desc, bg) in enumerate(col_data):
        cx = col_x[i]
        add_rect(slide, cx, 0, 3.1, 11.25, fill=bg)
        add_icon_circle(slide, sym, cx + 1.1, 1.2, 0.9, text_size=Pt(20))
        add_text(slide, title, cx + 0.2, 2.45, 2.7, 1.35,
                 font=FONT_DISPLAY, size=Pt(21), color=ORANGE)
        add_rule(slide, cx + 0.2, 4.05, 2.65)
        add_text(slide, desc, cx + 0.2, 4.25, 2.7, 6.5,
                 size=Pt(13), color=CREAM)

    add_footer(slide, 10)


def s11_ask(prs):
    """SLIDE 11 — Investment Ask"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BLACK)

    add_rect(slide, 0, 0, 20, 0.07, fill=ORANGE)
    add_label_tag(slide, "The Ask", 0.8, 0.35, 2.8)

    # Giant ask number
    add_text(slide, "Rs.20 LAKHS", 0.8, 0.9, 13, 2.3,
             font=FONT_DISPLAY, size=Pt(98), color=ORANGE)
    add_text(slide, "FOR 10% EQUITY", 0.8, 3.2, 13, 1.0,
             font=FONT_DISPLAY, size=Pt(44), color=WHITE)

    add_rule(slide, 0.8, 4.35, 18.5)
    add_text(slide, "FUND ALLOCATION", 0.8, 4.6, 12, 0.5,
             font=FONT_DISPLAY, size=Pt(18), color=CREAM)

    allocs = [
        ("PD", "Product Development & App Refinement",         40),
        ("MK", "Marketing & City Launch Campaigns",            30),
        ("PS", "Partnership Expansion — B2B and NGOs",         20),
        ("OP", "Operations & Customer Support Infrastructure", 10),
    ]
    bar_total = 10.5
    y = 5.3
    for sym, label, pct in allocs:
        add_icon_shape(slide, sym, 0.8, y, 0.5, text_size=Pt(11))
        add_text(slide, label, 1.55, y + 0.05, 6.8, 0.38,
                 size=Pt(14), color=CREAM)
        add_text(slide, f"{pct}%", 8.5, y + 0.05, 1.0, 0.38,
                 font=FONT_DISPLAY, size=Pt(16), color=ORANGE)
        # Track
        add_rect(slide, 0.8, y + 0.52, bar_total, 0.28, fill=MID_GRAY)
        # Fill
        add_rect(slide, 0.8, y + 0.52, bar_total * pct / 100, 0.28, fill=ORANGE)
        y += 1.06

    # Right vision panel
    add_rect(slide, 12.3, 0.85, 7.3, 9.2, fill=DARK_GRAY)
    add_rect(slide, 12.3, 0.85, 0.07, 9.2, fill=ORANGE)

    add_text(slide, "OUR VISION", 12.7, 1.15, 6.6, 0.6,
             font=FONT_DISPLAY, size=Pt(22), color=ORANGE)
    add_text(slide,
             "To become India's largest platform that transforms loneliness "
             "into meaningful human connection — one city, one community at a time.\n\n"
             "The technology scales. Only the communities grow.",
             12.7, 1.95, 6.6, 3.2,
             size=Pt(16), color=CREAM)

    add_rule(slide, 12.7, 5.4, 6.4)
    add_text(slide, "YEAR ONE TARGETS", 12.7, 5.6, 6.6, 0.52,
             font=FONT_DISPLAY, size=Pt(18), color=WHITE)

    targets = [
        ("20+",  "business partners onboarded"),
        ("15+",  "NGO integrations live"),
        ("10K+", "verified active users"),
        ("100+", "community events per month"),
    ]
    y2 = 6.3
    for num, label in targets:
        add_text(slide, num, 12.7, y2, 1.4, 0.46,
                 font=FONT_DISPLAY, size=Pt(18), color=ORANGE)
        add_text(slide, label, 14.25, y2 + 0.06, 5.0, 0.38,
                 size=Pt(14), color=LIGHT_GRAY)
        y2 += 0.6

    add_footer(slide, 11)


def s12_thanks(prs):
    """SLIDE 12 — Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, ORANGE)

    # ── Layout: LEFT black panel | CENTER orange zone | RIGHT black panel ──
    # Left black panel: 0 – 6.8 in
    add_rect(slide, 0,    0, 6.8,  11.25, fill=BLACK)
    # Right black panel: 13.2 – 20 in
    add_rect(slide, 13.2, 0, 6.8,  11.25, fill=BLACK)
    # Centre stays orange (background)

    # ── Centre: THANK YOU ──────────────────────────────────────────────────
    add_text(slide, "THANK", 6.8, 1.5, 6.4, 2.0,
             font=FONT_DISPLAY, size=Pt(110), color=BLACK,
             align=PP_ALIGN.CENTER)
    add_text(slide, "YOU.", 6.8, 3.6, 6.4, 2.0,
             font=FONT_DISPLAY, size=Pt(110), color=BLACK,
             align=PP_ALIGN.CENTER)
    add_text(slide, "For Your Attention", 7.2, 5.9, 5.6, 0.9,
             font=FONT_BODY, size=Pt(24), italic=True, color=BLACK,
             align=PP_ALIGN.CENTER)
    # Short rule under subtitle
    add_rule(slide, 9.0, 7.0, 2.0, thickness=0.06, color=BLACK)

    # ── Left panel: SocialSpark identity ──────────────────────────────────
    add_text(slide, "SOCIALSPARK", 0.45, 2.8, 6.0, 0.9,
             font=FONT_DISPLAY, size=Pt(30), color=ORANGE)
    add_rule(slide, 0.45, 3.85, 5.8, thickness=0.06, color=MID_GRAY)
    add_text(slide,
             '"Building Social Fitness,\nOne Connection at a Time."',
             0.45, 4.1, 6.0, 1.8,
             size=Pt(17), italic=True, color=CREAM)
    add_text(slide,
             '"Loneliness is not solved by more screen time.\nIt is solved by more real time."',
             0.45, 6.2, 6.0, 2.2,
             size=Pt(15), italic=True, color=LIGHT_GRAY)

    # ── Right panel: Contact ───────────────────────────────────────────────
    add_text(slide, "QUESTIONS?", 13.55, 2.8, 6.0, 0.9,
             font=FONT_DISPLAY, size=Pt(30), color=ORANGE)
    add_rule(slide, 13.55, 3.85, 5.8, thickness=0.06, color=MID_GRAY)
    add_text(slide, "Let's Connect.", 13.55, 4.15, 6.0, 0.6,
             size=Pt(20), color=WHITE)
    add_text(slide, "team@socialspark.in", 13.55, 4.85, 6.0, 0.55,
             size=Pt(17), bold=True, color=ORANGE)
    add_rule(slide, 13.55, 5.6, 5.8, thickness=0.04, color=MID_GRAY)
    add_text(slide,
             "Entrepreneurship Cell\nShark Tank Pitch Competition",
             13.55, 5.85, 6.0, 1.5,
             size=Pt(16), color=LIGHT_GRAY)

    # Footer stays on orange zone
    add_rect(slide, 0, 10.85, 20, 0.4, fill=BLACK)
    add_text(slide,
             "SocialSpark Pitch Deck  —  Shark Tank Competition  —  Slide 12 of 12",
             0.5, 10.9, 19, 0.3,
             size=Pt(11), color=ORANGE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────
# BUILD
# ─────────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    s01_cover(prs)
    s02_problem(prs)
    s03_gap(prs)
    s04_solution(prs)
    s05_howit(prs)
    s06_gamify(prs)
    s07_business(prs)
    s08_scale(prs)
    s09_usp(prs)
    s10_whynow(prs)
    s11_ask(prs)
    s12_thanks(prs)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")

if __name__ == "__main__":
    build()
